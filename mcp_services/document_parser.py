import asyncio
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional
import os
import json

# 导入文档处理库
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import PyPDF2
    import pdfplumber
except ImportError:
    PyPDF2 = None
    pdfplumber = None

from .base_service import BaseMCPService
from utils.file_utils import FileUtils

logger = logging.getLogger(__name__)

class DocumentParserService(BaseMCPService):
    """文档解析服务 - 简化版"""
    
    def __init__(self):
        super().__init__("document_parser")
        self.supported_formats = ['.docx', '.doc', '.pptx', '.xlsx', '.xls', '.pdf', '.txt']
    
    async def initialize(self) -> bool:
        """初始化服务"""
        # 注册方法处理器
        self.register_method("extract_text", self._extract_text)
        self.register_method("extract_metadata", self._extract_metadata)
        self.register_method("parse_document", self._parse_document)
        
        logger.info("文档解析服务初始化完成")
        return True

    async def _extract_text(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """提取文档文本内容"""
        file_path = params.get("file_path")
        if not file_path:
            raise ValueError("缺少file_path参数")
        
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        extension = file_path.suffix.lower()
        
        try:
            if extension in ['.txt']:
                text_content = await self._extract_txt_text(file_path)
            elif extension in ['.docx']:
                text_content = await self._extract_docx_text(file_path)
            elif extension in ['.pptx']:
                text_content = await self._extract_pptx_text(file_path)
            elif extension in ['.xlsx', '.xls']:
                text_content = await self._extract_excel_text(file_path)
            elif extension == '.pdf':
                text_content = await self._extract_pdf_text(file_path)
            else:
                raise ValueError(f"不支持的文件格式: {extension}")
            
            return {
                "text_content": text_content,
                "file_path": str(file_path),
                "file_type": extension,
                "content_length": len(text_content)
            }
            
        except Exception as e:
            logger.error(f"文本提取失败: {e}")
            raise

    async def _extract_metadata(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """提取文档元数据"""
        file_path = params.get("file_path")
        if not file_path:
            raise ValueError("缺少file_path参数")
        
        file_path = Path(file_path)
        file_info = FileUtils.get_file_info(file_path)
        
        return {
            "filename": file_info["name"],
            "size": file_info["size"],
            "size_human": file_info["size_human"],
            "extension": file_info["extension"],
            "mime_type": file_info["mime_type"],
            "created": file_info["created"],
            "modified": file_info["modified"]
        }

    async def _parse_document(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """解析文档（综合功能）"""
        file_path = params.get("file_path")
        extract_tables = params.get("extract_tables", True)
        extract_images = params.get("extract_images", False)
        
        # 提取文本
        text_result = await self._extract_text({"file_path": file_path})
        
        # 提取元数据
        metadata_result = await self._extract_metadata({"file_path": file_path})
        
        result = {
            "text_content": text_result["text_content"],
            "metadata": metadata_result,
            "tables": [],
            "images": []
        }
        
        # 根据需要提取表格和图片
        if extract_tables:
            try:
                tables = await self._extract_tables({"file_path": file_path})
                result["tables"] = tables.get("tables", [])
            except Exception as e:
                logger.warning(f"表格提取失败: {e}")
        
        if extract_images:
            try:
                images = await self._extract_images({"file_path": file_path})
                result["images"] = images.get("images", [])
            except Exception as e:
                logger.warning(f"图片提取失败: {e}")
        
        return result

    async def _extract_docx_text(self, file_path: Path) -> str:
        """提取DOCX文档文本"""
        if Document is None:
            raise ImportError("python-docx库未安装")
        
        doc = Document(file_path)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        return '\n'.join(text_parts)

    async def _extract_pptx_text(self, file_path: Path) -> str:
        """提取PPTX文档文本"""
        if Presentation is None:
            raise ImportError("python-pptx库未安装")
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return '\n'.join(text_parts)

    async def _extract_excel_text(self, file_path: Path) -> str:
        """提取Excel文档文本"""
        if openpyxl is None:
            raise ImportError("openpyxl库未安装")
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"工作表: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    text_parts.append('\t'.join(row_text))
        
        return '\n'.join(text_parts)

    async def _extract_pdf_text(self, file_path: Path) -> str:
        """提取PDF文档文本"""
        if pdfplumber is None:
            raise ImportError("pdfplumber库未安装")
        
        text_parts = []
        
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        
        return '\n'.join(text_parts)

    async def _extract_tables(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """提取文档中的表格"""
        file_path = Path(params.get("file_path"))
        extension = file_path.suffix.lower()
        
        tables = []
        
        if extension in ['.xlsx', '.xls'] and openpyxl:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                table_data = []
                for row in sheet.iter_rows(values_only=True):
                    table_data.append([str(cell) if cell is not None else "" for cell in row])
                
                if table_data:
                    tables.append({
                        "sheet_name": sheet_name,
                        "data": table_data,
                        "rows": len(table_data),
                        "columns": len(table_data[0]) if table_data else 0
                    })
        
        return {"tables": tables}

    async def _extract_images(self, params: Dict[str, Any]) -> Dict[str, Any]:
        """提取文档中的图片信息"""
        # 简化实现，返回空列表
        return {"images": []}

    async def _extract_txt_text(self, file_path: Path) -> str:
        """提取txt文件的文本内容"""
        try:
            # 尝试不同的编码格式
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                        logger.info(f"成功使用 {encoding} 编码读取文件")
                        return content
                except UnicodeDecodeError:
                    continue
            
            # 如果所有编码都失败，使用默认编码并忽略错误
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                logger.warning("使用UTF-8编码并忽略错误读取文件")
                return content
                
        except Exception as e:
            logger.error(f"读取txt文件失败: {e}")
            raise
        return {"images": []}
